"""PPTX 生成服務"""
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image
import io


class PptxService:
    """PPTX 生成服務類"""
    
    def __init__(self, ratio: str = "16:9"):
        """
        初始化簡報
        
        Args:
            ratio: 投影片比例 ("16:9" 或 "4:3")
        """
        self.prs = Presentation()
        
        # 設定投影片尺寸
        if ratio == "16:9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        else:  # 4:3
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
    
    def add_slide_with_background(self, background_image: Image.Image, texts: List[Dict]) -> None:
        """
        新增一張投影片，包含背景圖和文字
        
        Args:
            background_image: 背景圖片
            texts: 文字資料列表 [{content, x, y, width, height, font_size, color, ...}]
        """
        # 使用空白版面
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 儲存背景圖到 bytes
        img_bytes = io.BytesIO()
        background_image.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        # 添加背景圖（填滿整個投影片）
        slide.shapes.add_picture(
            img_bytes,
            Emu(0), Emu(0),
            self.prs.slide_width, self.prs.slide_height
        )
        
        # 計算縮放比例
        scale_x = self.prs.slide_width / Emu(background_image.width * 914400 / 96)
        scale_y = self.prs.slide_height / Emu(background_image.height * 914400 / 96)
        
        # 添加文字框
        for text_data in texts:
            self._add_text_box(slide, text_data, scale_x, scale_y)
    
    def _add_text_box(self, slide, text_data: Dict, scale_x: float, scale_y: float) -> None:
        """添加文字框到投影片"""
        # 驗證必要欄位
        content = text_data.get('content', '')
        if not content:
            return  # 跳過空內容
        
        # 計算位置（從像素轉換為 EMU）
        px_to_emu = 914400 / 96  # 96 DPI
        
        x = text_data.get('x', 0)
        y = text_data.get('y', 0)
        w = text_data.get('width', 200)
        h = text_data.get('height', 50)
        
        try:
            left = Emu(int(x * px_to_emu * scale_x))
            top = Emu(int(y * px_to_emu * scale_y))
            width = Emu(int(w * px_to_emu * scale_x))
            height = Emu(int(h * px_to_emu * scale_y))
            
            # 確保尺寸有效
            if width <= 0:
                width = Emu(int(200 * px_to_emu * scale_x))
            if height <= 0:
                height = Emu(int(50 * px_to_emu * scale_y))
            
            # 添加文字框
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = content
            
            # 設定字體大小
            font_size = text_data.get('font_size')
            if font_size and isinstance(font_size, (int, float)) and font_size > 0:
                p.font.size = Pt(int(font_size))
            
            # 設定粗體
            if text_data.get('font_weight') == 'bold':
                p.font.bold = True
            
            # 設定顏色
            color = text_data.get('color', '')
            if color and len(color) >= 6:
                hex_color = color.lstrip('#')
                if len(hex_color) >= 6:
                    try:
                        r = int(hex_color[0:2], 16)
                        g = int(hex_color[2:4], 16)
                        b = int(hex_color[4:6], 16)
                        p.font.color.rgb = RGBColor(r, g, b)
                    except ValueError:
                        pass  # 無效的顏色格式，跳過
        except Exception as e:
            print(f"Error adding text box: {e}")
            # 繼續處理其他文字框
    
    def save(self) -> bytes:
        """儲存並返回 PPTX bytes"""
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.read()
